#!/usr/bin/env python3
"""Generate the AgriSentinel Aaroh AI/ML Hackathon PowerPoint deck."""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
FRONTEND = ROOT / "frontend"
OUT = Path(__file__).resolve().parent / "AgriSentinel_Aaroh_AI_ML_Hackathon.pptx"

PRODUCT_SHOT = FRONTEND / "agrisentinel-product.png"
FIELD_IMAGE = ROOT / "agents/testdata/field_tomato_late_blight.jpg"
CONFUSION_MATRIX = ROOT / "ml/artifacts/confusion_matrix.png"


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


BG = rgb("07140F")
PANEL = rgb("0C2119")
PANEL_2 = rgb("102B21")
PANEL_3 = rgb("132F25")
WHITE = rgb("F5F8F6")
MUTED = rgb("A7B9B0")
DIM = rgb("6E8278")
GREEN = rgb("36D399")
GREEN_SOFT = rgb("75E6B8")
AMBER = rgb("F4C84A")
ORANGE = rgb("F97316")
RED = rgb("FF6B5E")
CYAN = rgb("5DD6D3")
BORDER = rgb("21473A")
CHARCOAL = rgb("0A1712")

W = 13.333
H = 7.5


prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
prs.core_properties.title = "AgriSentinel — Aaroh AI/ML Hackathon"
prs.core_properties.subject = "Hybrid crop disease detection, verified field action planning, and grounded follow-up advice"
prs.core_properties.author = "Vishal and Khushi"
prs.core_properties.keywords = "AgriSentinel, Aaroh, AI, ML, agriculture, EfficientNet, vision LLM, RAG, agents, advisor"


def set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color: RGBColor | None = None, width: float = 1.0) -> None:
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color
        shape.line.width = Pt(width)


def rect(slide, x, y, w, h, fill=PANEL, line=BORDER, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    set_line(shape, line, 0.8 if line else 0)
    return shape


def line(slide, x, y, w, h=0, color=BORDER, width=1.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(max(h, 0.01)))
    set_fill(shape, color)
    set_line(shape, None)
    return shape


def textbox(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=18,
    color=WHITE,
    bold=False,
    font="Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.03,
    fit=False,
    italic=False,
    tracking=None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    if fit:
        frame.fit_text(font_family=font, max_size=int(size))
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    paragraph.line_spacing = 1.05
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if tracking is not None:
        run.font._element.set("spc", str(tracking))
    return box


def bullets(slide, items, x, y, w, h, size=16, color=MUTED, gap=8, bullet_color=GREEN):
    item_h = h / max(len(items), 1)
    for idx, item in enumerate(items):
        yy = y + idx * item_h
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(yy + 0.12), Inches(0.09), Inches(0.09))
        set_fill(dot, bullet_color)
        set_line(dot, None)
        textbox(slide, item, x + 0.18, yy, w - 0.18, item_h, size=size, color=color, valign=MSO_ANCHOR.TOP)


def chip(slide, text, x, y, w, color=GREEN, fill=PANEL_2, size=10):
    shape = rect(slide, x, y, w, 0.34, fill=fill, line=color, radius=True)
    textbox(slide, text.upper(), x + 0.05, y + 0.01, w - 0.1, 0.30, size=size, color=color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return shape


def label(slide, text, x, y, w=2.5, color=GREEN):
    textbox(slide, text.upper(), x, y, w, 0.24, size=9, color=color, bold=True, tracking=120)


def logo(slide, x=0.55, y=0.34):
    outer = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(x), Inches(y), Inches(0.32), Inches(0.32))
    set_fill(outer, PANEL_3)
    set_line(outer, GREEN, 1.1)
    inner = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(x + 0.105), Inches(y + 0.105), Inches(0.11), Inches(0.11))
    set_fill(inner, GREEN_SOFT)
    set_line(inner, None)
    textbox(slide, "AGRISENTINEL", x + 0.46, y - 0.01, 2.1, 0.24, size=10, color=WHITE, bold=True)
    textbox(slide, "AUTONOMOUS FIELD HEALTH", x + 0.46, y + 0.17, 2.2, 0.16, size=6.5, color=GREEN, bold=True, tracking=90)


def base_slide(section: str, number: int, title: str | None = None, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    # restrained top-right grid motif
    for i in range(5):
        line(slide, 11.12 + i * 0.42, 0.0, 0.01, 0.75, color=BORDER, width=0)
    for i in range(3):
        line(slide, 11.1, 0.16 + i * 0.24, 2.23, 0.01, color=BORDER, width=0)
    logo(slide)
    textbox(slide, f"{number:02d}", 12.12, 0.31, 0.55, 0.25, size=9, color=DIM, bold=True, align=PP_ALIGN.RIGHT)
    textbox(slide, section.upper(), 10.05, 0.31, 1.95, 0.25, size=8, color=GREEN, bold=True, align=PP_ALIGN.RIGHT, tracking=100)
    if title:
        textbox(slide, title, 0.58, 0.98, 11.9, 0.65, size=28, color=WHITE, bold=True)
    if subtitle:
        textbox(slide, subtitle, 0.6, 1.62, 11.5, 0.42, size=13.5, color=MUTED)
    textbox(slide, "AAROH AI/ML HACKATHON  •  DEV A: VISHAL  •  DEV B: KHUSHI", 0.58, 7.16, 7.8, 0.18, size=7.5, color=DIM, bold=True, tracking=60)
    return slide


def add_notes(slide, text: str):
    frame = slide.notes_slide.notes_text_frame
    frame.text = text.strip()


def add_picture_cover(slide, path: Path, x, y, w, h, border=BORDER, border_w=1.0, focus_y=0.5):
    with Image.open(path) as image:
        iw, ih = image.size
    target = w / h
    source = iw / ih
    if source > target:
        crop = (1 - target / source) / 2
        crop_left = crop_right = crop
        crop_top = crop_bottom = 0
    else:
        visible = source / target
        crop_total = 1 - visible
        crop_top = max(0, min(crop_total, crop_total * focus_y))
        crop_bottom = crop_total - crop_top
        crop_left = crop_right = 0
    picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    picture.crop_left = crop_left
    picture.crop_right = crop_right
    picture.crop_top = crop_top
    picture.crop_bottom = crop_bottom
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    frame.fill.background()
    set_line(frame, border, border_w)
    return picture


def stat_card(slide, value, title, note, x, y, w=2.25, accent=GREEN):
    rect(slide, x, y, w, 1.38, fill=PANEL, line=BORDER)
    line(slide, x, y, 0.06, 1.38, color=accent)
    textbox(slide, value, x + 0.22, y + 0.20, w - 0.34, 0.48, size=25, color=WHITE, bold=True)
    textbox(slide, title.upper(), x + 0.22, y + 0.73, w - 0.34, 0.18, size=7.5, color=accent, bold=True, tracking=80)
    textbox(slide, note, x + 0.22, y + 0.98, w - 0.34, 0.22, size=8.5, color=DIM)


def flow_box(slide, title, subtitle, x, y, w, h, accent=GREEN, fill=PANEL, title_size=14):
    rect(slide, x, y, w, h, fill=fill, line=BORDER)
    line(slide, x, y, 0.06, h, color=accent)
    textbox(slide, title, x + 0.2, y + 0.16, w - 0.32, 0.32, size=title_size, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    textbox(slide, subtitle, x + 0.2, y + 0.53, w - 0.32, h - 0.66, size=9.5, color=MUTED)


def chevron(slide, x, y, w=0.3, h=0.32, color=GREEN):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, color)
    set_line(shape, None)
    return shape


# ---------------------------------------------------------------------------
# 1. Cover
# ---------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG
for idx in range(8):
    line(slide, 0, 0.6 + idx * 0.82, 13.333, 0.01, color=rgb("0B2118"))
for idx in range(14):
    line(slide, 0.55 + idx * 0.94, 0, 0.01, 7.5, color=rgb("0B2118"))
logo(slide, 0.65, 0.55)
chip(slide, "Aaroh AI/ML Hackathon", 0.67, 1.30, 2.42, color=AMBER, fill=rgb("282511"), size=8.5)
textbox(slide, "From image to\nverified field action.", 0.68, 1.92, 6.15, 1.58, size=34, color=WHITE, bold=True)
textbox(slide, "Hybrid crop disease detection that combines local EfficientNet speed with whole-image vision reasoning — then verifies every recommendation before it reaches a farmer.", 0.72, 3.67, 5.42, 1.15, size=15.5, color=MUTED)
chip(slide, "HYBRID DETECTION", 0.72, 5.18, 1.58, color=GREEN, fill=PANEL_2, size=7.5)
chip(slide, "RAG + VERIFIER", 2.44, 5.18, 1.44, color=CYAN, fill=PANEL_2, size=7.5)
chip(slide, "EN + HI BRIEF", 4.03, 5.18, 1.35, color=AMBER, fill=rgb("282511"), size=7.5)
textbox(slide, "TEAM OWNERSHIP", 0.72, 6.00, 1.55, 0.18, size=7.5, color=DIM, bold=True, tracking=100)
textbox(slide, "DEV A  ·  VISHAL", 0.72, 6.27, 2.00, 0.22, size=10, color=GREEN, bold=True, tracking=45)
textbox(slide, "Backend • ML • Agents", 0.72, 6.54, 2.25, 0.20, size=9.5, color=WHITE, bold=True)
textbox(slide, "DEV B  ·  KHUSHI", 3.35, 6.27, 2.05, 0.22, size=10, color=CYAN, bold=True, tracking=45)
textbox(slide, "Frontend • UI • Demo", 3.35, 6.54, 2.25, 0.20, size=9.5, color=WHITE, bold=True)

rect(slide, 7.10, 0.72, 5.63, 6.12, fill=CHARCOAL, line=BORDER)
add_picture_cover(slide, FIELD_IMAGE, 7.28, 0.90, 5.27, 3.30, border=BORDER, focus_y=0.5)
label(slide, "Prototype field mosaic", 7.48, 1.08, 2.1, color=AMBER)
textbox(slide, "40", 7.45, 4.56, 1.0, 0.55, size=30, color=WHITE, bold=True)
textbox(slide, "inspection tiles", 7.48, 5.05, 1.75, 0.25, size=9, color=MUTED)
textbox(slide, "+", 9.18, 4.62, 0.5, 0.4, size=23, color=DIM, bold=True, align=PP_ALIGN.CENTER)
textbox(slide, "2", 9.78, 4.56, 1.0, 0.55, size=30, color=WHITE, bold=True)
textbox(slide, "independent views", 9.78, 5.05, 1.75, 0.25, size=9, color=MUTED)
textbox(slide, "=", 11.25, 4.62, 0.5, 0.4, size=23, color=DIM, bold=True, align=PP_ALIGN.CENTER)
textbox(slide, "1", 11.80, 4.56, 0.6, 0.55, size=30, color=GREEN_SOFT, bold=True)
textbox(slide, "verified decision", 11.50, 5.05, 1.05, 0.36, size=9, color=GREEN_SOFT, align=PP_ALIGN.CENTER)
line(slide, 7.45, 5.55, 4.68, 0.01, color=BORDER)
textbox(slide, "LOCAL CNN", 7.46, 5.76, 1.25, 0.22, size=8, color=GREEN, bold=True, tracking=70)
textbox(slide, "maps where", 7.46, 6.06, 1.6, 0.24, size=12, color=WHITE, bold=True)
textbox(slide, "VISION LLM", 10.00, 5.76, 1.35, 0.22, size=8, color=CYAN, bold=True, tracking=70)
textbox(slide, "checks what", 10.00, 6.06, 1.65, 0.24, size=12, color=WHITE, bold=True)
add_notes(slide, """
Open with the reframe: AgriSentinel is not another leaf-classification app. It turns one field image into a verified decision. Introduce Vishal as Dev A for backend, ML, and agents, and Khushi as Dev B for frontend, UI, and the demo experience. Then preview the hybrid idea: the local CNN maps where disease may be, while an independent vision model checks what the whole image actually shows.
""")


# ---------------------------------------------------------------------------
# 2. Problem
# ---------------------------------------------------------------------------
slide = base_slide("Problem", 2, "The real problem is detection-to-action latency", "A disease label is useful only if it becomes a timely, targeted, and safe field decision.")

stages = [
    ("01", "Late detection", "Only a fraction of the field is inspected; infection can spread before it is noticed.", ORANGE),
    ("02", "Look-alike symptoms", "Blight, spots, stress, and deficiencies can be confused from a single close-up.", AMBER),
    ("03", "No spread intelligence", "A label does not reveal affected area, clusters, direction, or likely impact.", CYAN),
    ("04", "Unsafe next step", "Generic advice can invent a chemical or dosage without evidence or a veto.", RED),
]
for i, (num, title, body, accent) in enumerate(stages):
    x = 0.62 + i * 3.15
    rect(slide, x, 2.32, 2.80, 3.02, fill=PANEL, line=BORDER)
    textbox(slide, num, x + 0.22, 2.54, 0.55, 0.30, size=11, color=accent, bold=True)
    line(slide, x + 0.22, 2.97, 2.34, 0.04, color=accent)
    textbox(slide, title, x + 0.22, 3.25, 2.36, 0.54, size=17, color=WHITE, bold=True)
    textbox(slide, body, x + 0.22, 4.00, 2.34, 1.03, size=12, color=MUTED)
    if i < 3:
        chevron(slide, x + 2.90, 3.65, 0.22, 0.38, color=DIM)

rect(slide, 0.62, 5.72, 12.02, 0.92, fill=rgb("0E2B20"), line=GREEN)
textbox(slide, "OUR REFRAME", 0.90, 5.94, 1.20, 0.22, size=8, color=GREEN, bold=True, tracking=90)
textbox(slide, "Do not stop at “What disease is this?”  Answer “Where is it, how serious is it, and what can safely happen next?”", 2.23, 5.87, 9.93, 0.43, size=16, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
add_notes(slide, """
The problem is not simply identifying a leaf. Today the chain breaks at four places: symptoms are detected late, look-alike diseases cause uncertainty, no one measures spread, and generic advice can be unsafe. AgriSentinel is designed around the full detection-to-action delay.
""")


# ---------------------------------------------------------------------------
# 3. Solution
# ---------------------------------------------------------------------------
slide = base_slide("Solution", 3, "One upload. One autonomous path to action.", "The farmer sees the field evidence, the system’s reasoning, and a plan only after verification.")

steps = [
    ("1", "UPLOAD", "Drone or phone field image", GREEN),
    ("2", "INSPECT", "8 × 5 tile heatmap", ORANGE),
    ("3", "FUSE", "Local + whole-image views", CYAN),
    ("4", "VERIFY", "Grounding, dosage, allowlist", AMBER),
    ("5", "ACT", "Schedule + farmer brief", GREEN_SOFT),
]
for i, (n, head, body, accent) in enumerate(steps):
    x = 0.62 + i * 2.48
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.38), Inches(0.54), Inches(0.54))
    set_fill(circle, accent)
    set_line(circle, None)
    textbox(slide, n, x, 2.40, 0.54, 0.50, size=14, color=BG, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    if i < 4:
        line(slide, x + 0.54, 2.63, 1.92, 0.03, color=BORDER)
        chevron(slide, x + 2.21, 2.48, 0.25, 0.30, color=accent)
    textbox(slide, head, x, 3.10, 1.80, 0.24, size=9, color=accent, bold=True, tracking=90)
    textbox(slide, body, x, 3.48, 1.95, 0.62, size=13, color=WHITE, bold=True)

outcomes = [
    ("Infection map", "Progressive tile heatmap with confidence escalation", GREEN),
    ("Spread intelligence", "% affected, clusters, direction, yield-risk estimate", ORANGE),
    ("Verified plan", "Source-backed treatment with PASS / REWRITE / BLOCK", AMBER),
    ("Farmer action", "Day-by-day schedule, cost band, rescan, EN / HI brief", CYAN),
]
for i, (head, body, accent) in enumerate(outcomes):
    x = 0.62 + i * 3.03
    rect(slide, x, 4.66, 2.78, 1.50, fill=PANEL, line=BORDER)
    line(slide, x, 4.66, 2.78, 0.05, color=accent)
    textbox(slide, head, x + 0.20, 4.92, 2.35, 0.28, size=13, color=WHITE, bold=True)
    textbox(slide, body, x + 0.20, 5.36, 2.35, 0.55, size=9.7, color=MUTED)

textbox(slide, "The architecture changes control flow — not just wording.", 0.64, 6.48, 7.6, 0.30, size=13, color=GREEN_SOFT, bold=True)
add_notes(slide, """
Walk left to right. The user uploads once. The system tiles the image, fuses two independent visual reads, verifies the grounded plan, and produces an action schedule. The key product outcomes are visible: heatmap, spread, verified guidance, and a plain-language brief.
""")


# ---------------------------------------------------------------------------
# 4. Hybrid detection
# ---------------------------------------------------------------------------
slide = base_slide("Differentiator", 4, "Hybrid detection: local speed + global context", "EfficientNet-B0 and a vision LLM inspect the same image independently before a policy-driven fusion step.")

textbox(slide, "FIELD IMAGE", 0.66, 2.31, 1.15, 0.28, size=9, color=MUTED, bold=True, tracking=80)
flow_box(slide, "Image", "Drone / phone field capture", 0.64, 2.70, 1.55, 1.12, accent=WHITE, fill=PANEL)
chevron(slide, 2.33, 3.09, 0.35, 0.34, color=DIM)

flow_box(slide, "EfficientNet-B0", "Local ONNX model\n40 tile predictions\nMaps WHERE", 2.83, 2.18, 2.35, 1.70, accent=GREEN, fill=PANEL_2, title_size=15)
flow_box(slide, "Vision LLM", "Gemini / GPT-4.1\nWhole-image reasoning\nChecks WHAT", 2.83, 4.08, 2.35, 1.70, accent=CYAN, fill=PANEL_2, title_size=15)

chevron(slide, 5.42, 2.83, 0.35, 0.34, color=GREEN)
chevron(slide, 5.42, 4.73, 0.35, 0.34, color=CYAN)
flow_box(slide, "Decision Fusion Agent", "Reconciles crop, disease, visible symptoms, and disagreement policy", 5.90, 3.03, 2.67, 1.75, accent=AMBER, fill=rgb("282511"), title_size=16)

chevron(slide, 8.75, 3.73, 0.36, 0.34, color=AMBER)
flow_box(slide, "Final diagnosis", "Disease + confidence\nCNN spatial map preserved", 9.22, 3.03, 1.86, 1.75, accent=GREEN_SOFT, fill=PANEL_2, title_size=15)
chevron(slide, 11.23, 3.73, 0.34, 0.34, color=GREEN_SOFT)
flow_box(slide, "RAG → Verify → Plan", "Grounded treatment\nVeto before action", 11.68, 3.03, 1.02, 1.75, accent=RED, fill=PANEL, title_size=11)

rect(slide, 0.64, 6.24, 12.06, 0.55, fill=rgb("0E2B20"), line=GREEN)
textbox(slide, "Why hybrid?", 0.90, 6.39, 1.14, 0.20, size=9, color=GREEN, bold=True)
textbox(slide, "Tile-level precision without losing the whole-field context that can expose a confidently wrong local diagnosis.", 2.14, 6.32, 9.90, 0.28, size=12.5, color=WHITE, bold=True)
add_notes(slide, """
This is the new hybrid detection path. EfficientNet-B0 runs locally over forty tiles and is excellent at spatial localisation. A configurable vision LLM — Gemini today, or GPT-4.1 through the same role — reads the whole image without seeing the CNN answer. The Decision Fusion Agent reconciles them before RAG, verification, and planning.
""")


# ---------------------------------------------------------------------------
# 5. Fusion policy
# ---------------------------------------------------------------------------
slide = base_slide("Fusion logic", 5, "Disagreement becomes a safety signal", "The fusion layer preserves evidence and changes the outcome only when the independent views justify it.")

cards = [
    ("AGREE", "Same disease / healthy result", "Keep the CNN’s spatial heatmap and move forward with higher trust.", GREEN, "PASS FORWARD"),
    ("DISEASE MISMATCH", "Both see disease; names differ", "When vision is confident and crop-scoped, relabel infected tiles while preserving where the CNN found them.", CYAN, "FUSE THE NAME"),
    ("HEALTHY CONFLICT", "Vision sees clean; CNN sees ≥20% affected", "Keep the heatmap as evidence, mark the diagnosis contested, and withhold treatment advice.", RED, "REFUSE TO SPRAY"),
]
for i, (head, sub, body, accent, outcome) in enumerate(cards):
    x = 0.66 + i * 4.16
    rect(slide, x, 2.32, 3.82, 3.47, fill=PANEL, line=BORDER)
    line(slide, x, 2.32, 3.82, 0.07, color=accent)
    textbox(slide, head, x + 0.24, 2.63, 3.25, 0.24, size=9, color=accent, bold=True, tracking=85)
    textbox(slide, sub, x + 0.24, 3.08, 3.28, 0.56, size=17, color=WHITE, bold=True)
    textbox(slide, body, x + 0.24, 3.86, 3.26, 1.00, size=12, color=MUTED)
    chip(slide, outcome, x + 0.24, 5.13, 1.72, color=accent, fill=PANEL_2, size=7.5)

rect(slide, 0.66, 6.12, 12.02, 0.65, fill=PANEL_2, line=BORDER)
textbox(slide, "Graceful fallback", 0.91, 6.31, 1.52, 0.23, size=10, color=AMBER, bold=True)
textbox(slide, "If the vision API is unavailable, the event is logged and the local offline pipeline continues — no crash, no hidden state.", 2.55, 6.24, 9.58, 0.33, size=12.5, color=WHITE, bold=True)
add_notes(slide, """
The fusion policy has three important branches. Agreement passes forward. A disease-name mismatch can fuse the label while keeping the CNN’s tile map. But a clean-versus-infected conflict is fundamentally different: the system marks the result contested and refuses to recommend spraying. If the vision API is unavailable, the local pipeline continues and the event log makes that visible.
""")


# ---------------------------------------------------------------------------
# 6. End-to-end architecture
# ---------------------------------------------------------------------------
slide = base_slide(
    "System architecture",
    6,
    "Architecture: pixels to a verified decision",
    "The browser, API, agents, models, knowledge, and safety boundary stay decoupled through one shared RunState.",
)

# Primary request path.
flow_box(slide, "Field capture", "Drone / phone\nJPEG or PNG", 0.62, 2.52, 1.43, 1.24, accent=WHITE, fill=PANEL, title_size=13)
chevron(slide, 2.15, 2.96, 0.25, 0.31, color=DIM)
flow_box(slide, "React console", "Upload • heatmap\nagent timeline\nbrief + Advisor", 2.49, 2.36, 1.79, 1.56, accent=GREEN, fill=PANEL_2, title_size=14)
chevron(slide, 4.39, 2.96, 0.25, 0.31, color=GREEN)
flow_box(slide, "FastAPI", "POST run • GET state\nSSE events • chat", 4.73, 2.36, 1.76, 1.56, accent=CYAN, fill=PANEL_2, title_size=14)
chevron(slide, 6.60, 2.96, 0.25, 0.31, color=CYAN)

# Supervised intelligence boundary.
rect(slide, 6.96, 2.16, 5.72, 3.77, fill=CHARCOAL, line=AMBER)
label(slide, "Supervised intelligence", 7.25, 2.42, 2.35, color=AMBER)
flow_box(slide, "Orchestrator", "Owns order, gates,\nbranches, and the\nverify-redraft loop", 7.24, 2.86, 1.38, 2.18, accent=AMBER, fill=rgb("282511"), title_size=13)

pipeline_groups = [
    ("PERCEPTION", "Scout • Observer\nDiagnose • 2nd opinion", GREEN),
    ("FUSION", "Consensus • Spread\nmap + severity", CYAN),
    ("GROUND", "TF-IDF RAG\n10 curated documents", AMBER),
    ("SAFETY", "Verifier → PASS\nREWRITE / BLOCK", RED),
    ("ACTION", "Planner • Reporter\nEN / HI brief", GREEN_SOFT),
]
for i, (head, body, accent) in enumerate(pipeline_groups):
    x = 8.88 + (i % 3) * 1.15
    y = 2.83 + (i // 3) * 1.20
    w = 1.02 if i < 3 else 1.60
    if i == 4:
        x = 10.68
    rect(slide, x, y, w, 0.98, fill=PANEL, line=accent)
    textbox(slide, head, x + 0.09, y + 0.11, w - 0.18, 0.18, size=6.7, color=accent, bold=True, tracking=35, align=PP_ALIGN.CENTER)
    textbox(slide, body, x + 0.08, y + 0.38, w - 0.16, 0.44, size=7.4, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    if i in (0, 1, 3):
        chevron(slide, x + w + 0.04, y + 0.36, 0.14, 0.24, color=accent)

# State and service foundations.
rect(slide, 2.49, 4.34, 3.99, 1.58, fill=PANEL, line=BORDER)
label(slide, "Shared state + realtime", 2.74, 4.60, 2.15, color=GREEN)
textbox(slide, "Immutable contract-shaped RunState", 2.74, 4.96, 3.30, 0.25, size=12.5, color=WHITE, bold=True)
textbox(slide, "events[] is replayed over SSE; SQLite + uploads persist completed runs", 2.74, 5.31, 3.32, 0.39, size=9.3, color=MUTED)
line(slide, 3.34, 4.17, 0.03, 0.17, color=GREEN)
textbox(slide, "READ / WRITE", 4.91, 4.10, 1.14, 0.18, size=6.5, color=DIM, bold=True, align=PP_ALIGN.RIGHT, tracking=40)

rect(slide, 7.24, 5.30, 5.16, 0.38, fill=PANEL_2, line=BORDER)
textbox(slide, "ONNX Runtime  •  Vision LLM  •  TF-IDF retrieval  •  allowlist + dose rules", 7.41, 5.39, 4.82, 0.17, size=7.6, color=MUTED, bold=True, align=PP_ALIGN.CENTER)

# Explicit return paths keep the architecture legible in a five-minute pitch.
line(slide, 0.96, 4.23, 1.39, 0.03, color=GREEN)
chevron(slide, 0.82, 4.09, 0.18, 0.29, color=GREEN)
textbox(slide, "SSE EVENTS + FINAL STATE", 0.62, 4.46, 1.66, 0.35, size=7.2, color=GREEN_SOFT, bold=True, align=PP_ALIGN.CENTER)
textbox(slide, "Live scan", 0.62, 5.08, 0.80, 0.20, size=8, color=GREEN, bold=True)
textbox(slide, "Offline replay", 0.62, 5.40, 0.98, 0.20, size=8, color=CYAN, bold=True)
textbox(slide, "Same event consumer", 0.62, 5.70, 1.49, 0.20, size=7.8, color=MUTED)

rect(slide, 0.62, 6.23, 12.06, 0.53, fill=rgb("0E2B20"), line=GREEN)
textbox(slide, "CONTROL PRINCIPLE", 0.88, 6.39, 1.47, 0.20, size=7.5, color=GREEN, bold=True, tracking=75)
textbox(slide, "Agents never call each other; they read and write RunState, while the Orchestrator alone decides what runs next.", 2.47, 6.31, 9.73, 0.28, size=11.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_notes(slide, """
Read this from left to right. The React console uploads the field image to FastAPI, which starts a background run and returns the run ID immediately. The Orchestrator is the only component that controls sequence, confidence escalation, and the verifier rewrite loop. Every specialised agent reads and writes the shared RunState; the browser receives the same event history over SSE. ONNX inference remains local, while the whole-image Observer and optional language-model work can use a hosted provider. SQLite and uploaded images make finished runs recoverable. The offline demo replays the same events into the same frontend consumer.
""")


# ---------------------------------------------------------------------------
# 7. Agent system
# ---------------------------------------------------------------------------
slide = base_slide("Agent system", 7, "Eleven visible agents, one controlled execution trace", "Ten specialised agents share one event-sourced RunState; the Orchestrator alone controls sequence and branching.")

groups = [
    ("PERCEPTION", GREEN, [
        ("Scout", "tile + filter"),
        ("Observer", "whole-image vision"),
        ("Diagnostician", "EfficientNet-B0"),
        ("Second Opinion", "TTA under 0.75"),
        ("Consensus", "decision fusion"),
    ]),
    ("FIELD INTELLIGENCE", CYAN, [
        ("Spread Analyst", "% + clusters + direction"),
        ("RAG Agronomist", "source-grounded draft"),
    ]),
    ("SAFETY + ACTION", AMBER, [
        ("Verifier", "PASS / REWRITE / BLOCK"),
        ("Action Planner", "schedule + cost"),
        ("Reporter", "farmer brief EN / HI"),
    ]),
]

positions = [(0.64, 2.28, 4.12), (4.91, 2.28, 3.22), (8.29, 2.28, 4.40)]
for (group_name, accent, agents), (gx, gy, gw) in zip(groups, positions):
    rect(slide, gx, gy, gw, 3.83, fill=PANEL, line=BORDER)
    label(slide, group_name, gx + 0.22, gy + 0.22, gw - 0.44, color=accent)
    available = 3.15
    step = available / len(agents)
    for j, (name, role) in enumerate(agents):
        yy = gy + 0.70 + j * step
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(gx + 0.22), Inches(yy + 0.03), Inches(0.24), Inches(0.24))
        set_fill(dot, accent)
        set_line(dot, None)
        textbox(slide, str(j + 1), gx + 0.22, yy + 0.025, 0.24, 0.24, size=7, color=BG, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        textbox(slide, name, gx + 0.58, yy, gw * 0.43, 0.27, size=11.5, color=WHITE, bold=True)
        textbox(slide, role, gx + gw * 0.52, yy + 0.01, gw * 0.40, 0.24, size=8.8, color=MUTED, align=PP_ALIGN.RIGHT)

textbox(slide, "ORCHESTRATOR", 0.66, 6.42, 1.36, 0.20, size=8, color=GREEN, bold=True, tracking=80)
line(slide, 2.18, 6.50, 7.05, 0.03, color=GREEN)
textbox(slide, "Shared RunState + SSE events → live agent timeline", 9.43, 6.33, 3.20, 0.37, size=11, color=GREEN_SOFT, bold=True, align=PP_ALIGN.RIGHT)
add_notes(slide, """
There are eleven visible nodes including the Orchestrator. The multi-agent design is justified by different tools and different failure modes. Perception agents operate on pixels, the spread agent performs deterministic spatial analysis, the agronomist retrieves sources, and the verifier has veto power. All mutations and branches appear in one event stream, so judges can see what ran instead of trusting an architecture claim.
""")


# ---------------------------------------------------------------------------
# 8. Product
# ---------------------------------------------------------------------------
slide = base_slide("Product", 8, "The full decision is visible in one console", "This screenshot is the running React interface in offline replay mode, using the same event consumer as the live API.")

rect(slide, 0.60, 2.17, 8.72, 4.64, fill=CHARCOAL, line=BORDER)
add_picture_cover(slide, PRODUCT_SHOT, 0.72, 2.29, 8.48, 4.40, border=BORDER, focus_y=0.0)

callouts = [
    ("18.4%", "field affected", ORANGE),
    ("3", "infection clusters", AMBER),
    ("NE", "spread direction", CYAN),
    ("7/7", "claims grounded", GREEN),
]
for i, (value, desc, accent) in enumerate(callouts):
    y = 2.29 + i * 0.97
    rect(slide, 9.60, y, 3.06, 0.78, fill=PANEL, line=BORDER)
    textbox(slide, value, 9.83, y + 0.13, 0.95, 0.38, size=20, color=accent, bold=True)
    textbox(slide, desc.upper(), 10.87, y + 0.23, 1.55, 0.20, size=7.5, color=MUTED, bold=True, tracking=55)

rect(slide, 9.60, 6.20, 3.06, 0.49, fill=rgb("0E2B20"), line=GREEN)
textbox(slide, "Same UI: live SSE or offline replay", 9.78, 6.31, 2.68, 0.18, size=9, color=GREEN_SOFT, bold=True, align=PP_ALIGN.CENTER)
add_notes(slide, """
This is the actual dashboard. Point out the progressive heatmap, spread cards, live agent pipeline, grounded treatment plan, and action schedule. The numbers shown are from the prepared light-infection demo case. Offline replay is deliberate demo insurance: it feeds the same frontend event consumer as the live SSE backend.
""")


# ---------------------------------------------------------------------------
# 9. Advisor chatbot
# ---------------------------------------------------------------------------
slide = base_slide(
    "Advisor chatbot",
    9,
    "A run-grounded Advisor — not a generic chatbot",
    "A stateless, run-scoped Advisor sits inside the Farmer Brief and inherits the Verifier’s safety decision.",
)

# Run-scoped answer path.
flow_box(slide, "1  Ask", "Question + recent\nbrowser transcript", 0.64, 2.33, 1.58, 1.20, accent=GREEN, fill=PANEL_2, title_size=13)
chevron(slide, 2.30, 2.76, 0.25, 0.31, color=GREEN)
flow_box(slide, "2  Bind", "Load completed\nRunState read-only", 2.64, 2.33, 1.58, 1.20, accent=CYAN, fill=PANEL_2, title_size=13)
chevron(slide, 4.30, 2.76, 0.25, 0.31, color=CYAN)
flow_box(slide, "3  Retrieve", "Disease-scoped +\ngeneral practice", 4.64, 2.33, 1.58, 1.20, accent=AMBER, fill=rgb("282511"), title_size=13)
chevron(slide, 6.30, 2.76, 0.25, 0.31, color=AMBER)
flow_box(slide, "4  Answer", "LLM when online;\nextractive fallback", 6.64, 2.33, 1.58, 1.20, accent=CYAN, fill=PANEL_2, title_size=13)

flow_box(slide, "5  Verify each sentence", "Recompute citation • drop unsupported text • reject unsafe chemical or dose", 2.10, 4.08, 2.86, 1.27, accent=RED, fill=PANEL, title_size=14)
chevron(slide, 5.10, 4.53, 0.26, 0.32, color=RED)
flow_box(slide, "6  Return a typed outcome", "answer • sources • grounded • refused • provider", 5.51, 4.08, 2.70, 1.27, accent=GREEN_SOFT, fill=PANEL, title_size=14)

textbox(slide, "POST  /api/run/{run_id}/chat", 0.67, 3.76, 1.50, 0.38, size=8.2, color=DIM, bold=True, align=PP_ALIGN.CENTER)
textbox(slide, "No chat state is written back to the frozen run contract.", 0.74, 4.32, 1.26, 0.72, size=9.3, color=MUTED, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

# Product mock-up: brief with embedded cited follow-up.
rect(slide, 8.60, 2.27, 4.07, 3.76, fill=CHARCOAL, line=BORDER)
label(slide, "Inside the farmer brief", 8.88, 2.54, 2.25, color=GREEN)
textbox(slide, "ASK ABOUT THIS FIELD", 8.88, 2.91, 2.16, 0.19, size=7.2, color=MUTED, bold=True, tracking=65)
rect(slide, 9.72, 3.24, 2.56, 0.54, fill=rgb("103126"), line=GREEN)
textbox(slide, "Can I spray before rain?", 9.88, 3.39, 2.22, 0.20, size=9.5, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
rect(slide, 8.88, 3.98, 3.15, 1.02, fill=PANEL, line=BORDER)
textbox(slide, "Avoid spraying if rain is expected within four hours; the application may wash off.", 9.05, 4.14, 2.79, 0.47, size=9.1, color=WHITE)
chip(slide, "1  SOURCE • P2", 9.05, 4.65, 1.24, color=GREEN, fill=PANEL_2, size=6.5)
rect(slide, 8.88, 5.23, 3.50, 0.49, fill=rgb("282511"), line=AMBER)
textbox(slide, "BLOCK / contested → advice withheld", 9.04, 5.36, 3.18, 0.20, size=8.7, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

rect(slide, 0.64, 6.27, 12.03, 0.50, fill=rgb("0E2B20"), line=GREEN)
textbox(slide, "SAFETY INHERITANCE", 0.91, 6.42, 1.63, 0.19, size=7.3, color=GREEN, bold=True, tracking=70)
textbox(slide, "If the scan was BLOCKED or the models disagree, the Advisor never retrieves or calls a model—it explains the refusal and directs the farmer to an extension officer.", 2.67, 6.34, 9.46, 0.27, size=10.4, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_notes(slide, """
The Advisor appears inside the Farmer Brief after a completed live run. The browser posts the question and a bounded recent transcript to a run-scoped endpoint. The server loads that RunState read-only, retrieves disease-specific and crop-agnostic passages from the same ten-document corpus, and asks the language model only when available. Every returned sentence is attributed again against the evidence, and the same chemical allowlist and dosage checks used by the Verifier are applied. If the model is unavailable, the closest extractive passage is returned with honest framing. If the run was blocked or contested, the Advisor never reaches retrieval or the network; it explains why advice was withheld.
""")


# ---------------------------------------------------------------------------
# 10. Metrics
# ---------------------------------------------------------------------------
slide = base_slide("Evidence", 10, "Measured performance, not presentation claims", "Results below come from committed artifacts and a 2,399-image, 14-class PlantVillage test split.")

stat_card(slide, "95.6%", "test accuracy", "clean test split", 0.62, 2.25, 2.24, GREEN)
stat_card(slide, "94.7%", "macro-F1", "+2.6 pp vs scratch", 3.00, 2.25, 2.24, CYAN)
stat_card(slide, "99.3%", "confident accuracy", "confidence ≥ 0.75", 0.62, 3.83, 2.24, AMBER)
stat_card(slide, "6.5 ms", "CPU / tile", "ONNX Runtime mean", 3.00, 3.83, 2.24, ORANGE)
stat_card(slide, "0.46 s", "offline scan", "median end to end", 0.62, 5.41, 2.24, GREEN_SOFT)
stat_card(slide, "91.1%", "mosaic tile accuracy", "labelled synthetic grid", 3.00, 5.41, 2.24, RED)

rect(slide, 5.56, 2.25, 7.10, 4.54, fill=rgb("F4F7F5"), line=BORDER)
slide.shapes.add_picture(str(CONFUSION_MATRIX), Inches(5.78), Inches(2.43), Inches(6.66), Inches(4.08))
textbox(slide, "EfficientNet-B0 confusion matrix • 14 classes", 7.36, 6.49, 3.48, 0.18, size=8, color=rgb("43564D"), bold=True, align=PP_ALIGN.CENTER)
textbox(slide, "Sources: ml/artifacts/metrics.json, latency.md, pipeline_scan.md, baseline_comparison.md", 0.63, 6.90, 11.92, 0.17, size=7.2, color=DIM)
add_notes(slide, """
State the evidence boundary before the numbers: these are PlantVillage test-split and labelled synthetic-mosaic results, not real-field accuracy. The fine-tuned EfficientNet reaches 95.6 percent accuracy and 94.7 percent macro-F1. ONNX CPU inference averages 6.5 milliseconds per tile, and the entire offline pipeline has a 0.46-second median.
""")


# ---------------------------------------------------------------------------
# 11. Honesty / domain gap
# ---------------------------------------------------------------------------
slide = base_slide("Validation", 11, "The lab-to-field gap is real — so we expose it", "Image degradation alone costs 9.0 accuracy points and more than doubles the escalation load.")

chart_x = 0.74
chart_y = 2.44
chart_w = 7.05
chart_h = 3.62
rect(slide, chart_x, chart_y, chart_w, chart_h, fill=PANEL, line=BORDER)
labels_data = [
    ("Clean test", 95.6, GREEN),
    ("Mild simulation", 92.8, CYAN),
    ("Field simulation", 86.6, ORANGE),
]
for idx, (name, value, accent) in enumerate(labels_data):
    yy = chart_y + 0.62 + idx * 0.90
    textbox(slide, name, chart_x + 0.30, yy - 0.02, 1.35, 0.26, size=10.5, color=MUTED, bold=True)
    rect(slide, chart_x + 1.72, yy, 4.43, 0.32, fill=rgb("12261E"), line=None, radius=True)
    width = 4.43 * value / 100.0
    rect(slide, chart_x + 1.72, yy, width, 0.32, fill=accent, line=None, radius=True)
    textbox(slide, f"{value:.1f}%", chart_x + 6.19, yy - 0.04, 0.55, 0.29, size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
line(slide, chart_x + 1.72, chart_y + 3.28, 4.43, 0.02, color=BORDER)
textbox(slide, "Accuracy under progressively harder image conditions", chart_x + 1.74, chart_y + 3.36, 4.70, 0.20, size=8, color=DIM)

rect(slide, 8.06, 2.44, 4.60, 1.65, fill=PANEL, line=BORDER)
textbox(slide, "15.8%  →  35.9%", 8.36, 2.78, 3.85, 0.40, size=24, color=AMBER, bold=True)
textbox(slide, "tiles escalated below 0.75 confidence", 8.38, 3.30, 3.63, 0.26, size=10.5, color=MUTED)

rect(slide, 8.06, 4.34, 4.60, 1.72, fill=rgb("282511"), line=AMBER)
label(slide, "Evidence boundary", 8.36, 4.62, 2.20, color=AMBER)
textbox(slide, "Original real-field benchmark: not collected yet", 8.36, 5.02, 3.74, 0.35, size=14, color=WHITE, bold=True)
textbox(slide, "Next: 30–50 untouched phone photos with provenance, then report the measured gap even if it is large.", 8.36, 5.47, 3.78, 0.43, size=9.5, color=MUTED)

textbox(slide, "Source: ml/artifacts/lab_vs_field.json • simulated “field” changes quality, not scene composition", 0.74, 6.37, 11.86, 0.22, size=7.4, color=DIM)
add_notes(slide, """
This is the honesty slide. Under simulated field degradation, accuracy falls from 95.6 to 86.6 percent and escalation rises from 15.8 to 35.9 percent. That simulation still does not add real clutter, overlapping leaves, soil, or camera variation, so it is not a substitute for original field photos. The next validation step is explicit.
""")


# ---------------------------------------------------------------------------
# 12. Safety
# ---------------------------------------------------------------------------
slide = base_slide("Safety", 12, "Refusal is a successful product outcome", "The Verifier can revise or stop treatment advice while still returning the field evidence and diagnosis context.")

flow_box(slide, "RAG treatment draft", "Every sentence carries a retrieved source marker", 0.68, 2.48, 2.55, 1.50, accent=CYAN, fill=PANEL)
chevron(slide, 3.41, 3.03, 0.35, 0.34, color=CYAN)
flow_box(slide, "Verifier", "Grounding • allowlist • dosage sanity", 3.93, 2.48, 2.36, 1.50, accent=AMBER, fill=rgb("282511"))

paths = [
    ("PASS", "Show plan", GREEN, 6.88, 2.18),
    ("REWRITE", "Regenerate / strip", AMBER, 6.88, 3.35),
    ("BLOCK", "Withhold treatment", RED, 6.88, 4.52),
]
for head, body, accent, x, y in paths:
    chevron(slide, 6.46, y + 0.34, 0.27, 0.29, color=accent)
    rect(slide, x, y, 2.05, 0.90, fill=PANEL, line=accent)
    textbox(slide, head, x + 0.18, y + 0.16, 0.83, 0.25, size=11, color=accent, bold=True)
    textbox(slide, body, x + 0.18, y + 0.49, 1.65, 0.20, size=8.5, color=MUTED)

rect(slide, 9.33, 2.18, 3.35, 3.24, fill=PANEL, line=BORDER)
label(slide, "Adversarial test", 9.60, 2.47, 1.90, color=RED)
textbox(slide, "10", 9.60, 2.89, 0.86, 0.52, size=30, color=WHITE, bold=True)
textbox(slide, "crafted attacks", 10.37, 3.06, 1.50, 0.24, size=10, color=MUTED)
line(slide, 9.60, 3.62, 2.62, 0.02, color=BORDER)
textbox(slide, "2 landed", 9.60, 3.86, 1.26, 0.31, size=17, color=AMBER, bold=True)
textbox(slide, "2 caught", 10.95, 3.86, 1.25, 0.31, size=17, color=GREEN, bold=True)
textbox(slide, "0 leaked to the farmer", 9.60, 4.43, 2.60, 0.31, size=15, color=WHITE, bold=True)
textbox(slide, "Small-sample counts — not a claimed universal rate.", 9.60, 4.93, 2.52, 0.25, size=8.5, color=DIM)

rect(slide, 0.68, 5.87, 12.0, 0.69, fill=rgb("2A1715"), line=RED)
textbox(slide, "BLOCK means no dosage, schedule, cost, or rescan prescription — only the evidence, reason, and a safe referral.", 0.96, 6.07, 11.35, 0.28, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
textbox(slide, "Source: ml/artifacts/block_rate.md • 10 attacks, 2 reached the plan, both were blocked", 0.68, 6.76, 11.7, 0.18, size=7.4, color=DIM)
add_notes(slide, """
The verifier is a real control boundary. It checks whether claims are supported, whether chemicals are allowlisted, and whether dosage formats and ranges are sane. In the prepared adversarial test, ten attacks were attempted, two reached the plan, both were caught, and zero leaked. Call these counts, not a universal safety rate.
""")


# ---------------------------------------------------------------------------
# 13. Team + feasibility
# ---------------------------------------------------------------------------
slide = base_slide("Execution", 13, "Two developers. One contract. One integrated demo.", "Dev A and Dev B worked in parallel against the same RunState and endpoint contract.")

rect(slide, 0.66, 2.32, 3.72, 3.34, fill=PANEL, line=BORDER)
chip(slide, "DEV A  •  VISHAL", 0.91, 2.61, 1.56, color=GREEN, fill=PANEL_2, size=8.5)
textbox(slide, "AI Systems + Backend", 0.91, 3.11, 2.74, 0.34, size=17.5, color=WHITE, bold=True)
bullets(slide, ["Model training + ONNX inference", "Orchestration, agents, RAG, safety", "FastAPI, SSE, storage, evaluation"], 0.92, 3.68, 3.00, 1.43, size=10.2, color=MUTED, gap=3, bullet_color=GREEN)

rect(slide, 4.57, 2.32, 3.72, 3.34, fill=PANEL, line=BORDER)
chip(slide, "DEV B  •  KHUSHI", 4.82, 2.61, 1.58, color=CYAN, fill=PANEL_2, size=8.5)
textbox(slide, "Product UI + Demo", 4.82, 3.11, 2.65, 0.34, size=17.5, color=WHITE, bold=True)
bullets(slide, ["React console + responsive field UX", "Heatmap, agent timeline, safety UI", "Offline replay, brief, Advisor experience"], 4.83, 3.68, 3.00, 1.43, size=10.2, color=MUTED, gap=3, bullet_color=CYAN)

rect(slide, 8.48, 2.32, 4.19, 3.34, fill=PANEL, line=BORDER)
label(slide, "Shared integration boundary", 8.75, 2.64, 2.95, color=AMBER)
stack = [
    ("CONTRACT", "Frozen RunState JSON schema", GREEN),
    ("TRANSPORT", "POST / GET / SSE / health / chat", CYAN),
    ("DEV A", "Publishes state + ordered events", AMBER),
    ("DEV B", "Renders live or replay identically", ORANGE),
]
for i, (head, body, accent) in enumerate(stack):
    yy = 3.18 + i * 0.58
    textbox(slide, head, 8.75, yy, 1.05, 0.21, size=7.5, color=accent, bold=True, tracking=50)
    textbox(slide, body, 9.93, yy - 0.02, 2.40, 0.24, size=9.8, color=WHITE, bold=True)
    if i < 3:
        line(slide, 8.75, yy + 0.36, 3.48, 0.01, color=BORDER)

rect(slide, 0.66, 6.10, 12.01, 0.58, fill=rgb("0E2B20"), line=GREEN)
textbox(slide, "Integration rule", 0.93, 6.28, 1.28, 0.22, size=9, color=GREEN, bold=True)
textbox(slide, "Dev A owns the system behind the contract; Dev B owns every user interaction in front of it.", 2.33, 6.20, 9.75, 0.31, size=11.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_notes(slide, """
Vishal is Dev A and owns the backend, model pipeline, orchestration, agents, retrieval, verifier, Advisor, and evaluation evidence. Khushi is Dev B and owns the React product experience: upload, heatmap, activity timeline, safety states, farmer brief, Advisor UI, responsive behaviour, and demo replay. The frozen RunState and endpoint contract let both developers work in parallel and made final integration predictable.
""")


# ---------------------------------------------------------------------------
# 14. Close
# ---------------------------------------------------------------------------
slide = base_slide("Vision", 14)
textbox(slide, "AgriSentinel turns uncertainty\ninto a safer field decision.", 0.70, 1.15, 7.62, 1.33, size=31, color=WHITE, bold=True)
textbox(slide, "Not just “what disease?” — but where it is, how serious it may be, and whether the next action is evidence-backed enough to show.", 0.74, 2.76, 6.82, 0.88, size=15, color=MUTED)

takeaways = [
    ("SEE", "Hybrid detection combines local spatial precision with whole-image context.", GREEN),
    ("VERIFY", "RAG grounding and a hard veto stop unsupported treatment from shipping.", AMBER),
    ("ACT", "The farmer gets a field map, schedule, cost band, and EN / HI brief.", CYAN),
]
for i, (head, body, accent) in enumerate(takeaways):
    y = 4.04 + i * 0.88
    chip(slide, head, 0.74, y, 0.92, color=accent, fill=PANEL_2, size=8.5)
    textbox(slide, body, 1.90, y - 0.01, 5.68, 0.38, size=12.5, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)

rect(slide, 8.35, 1.25, 4.32, 5.36, fill=PANEL, line=BORDER)
label(slide, "Prototype scope", 8.69, 1.65, 2.0, color=GREEN)
textbox(slide, "Tomato  •  Potato  •  Corn", 8.69, 2.07, 3.42, 0.34, size=17, color=WHITE, bold=True)
textbox(slide, "14 disease / healthy classes", 8.69, 2.52, 3.18, 0.24, size=10.5, color=MUTED)
line(slide, 8.69, 2.98, 3.46, 0.02, color=BORDER)
label(slide, "Next validation", 8.69, 3.30, 2.1, color=AMBER)
textbox(slide, "Real-field benchmark", 8.69, 3.72, 3.10, 0.34, size=17, color=WHITE, bold=True)
textbox(slide, "30–50 original photos • crop/device calibration • regional agronomy review", 8.69, 4.17, 3.20, 0.70, size=10.5, color=MUTED)
line(slide, 8.69, 5.02, 3.46, 0.02, color=BORDER)
textbox(slide, "THANK YOU", 8.69, 5.40, 1.42, 0.22, size=8, color=GREEN, bold=True, tracking=100)
textbox(slide, "Dev A: Vishal", 8.69, 5.69, 2.85, 0.28, size=14, color=WHITE, bold=True)
textbox(slide, "Dev B: Khushi", 8.69, 5.99, 2.85, 0.28, size=14, color=WHITE, bold=True)
textbox(slide, "Aaroh AI/ML Hackathon", 8.69, 6.31, 2.80, 0.22, size=9.5, color=GREEN_SOFT)
add_notes(slide, """
Close with three words: see, verify, act. AgriSentinel combines two visual perspectives, turns evidence into an auditable decision, and knows when to refuse. The immediate next step is not more features; it is the real-field benchmark and calibration needed to earn deployment trust. Thank the judges.
""")


def validate_assets():
    missing = [path for path in (PRODUCT_SHOT, FIELD_IMAGE, CONFUSION_MATRIX) if not path.exists()]
    if missing:
        raise FileNotFoundError("Missing presentation assets: " + ", ".join(map(str, missing)))


def validate_deck():
    assert len(prs.slides) == 14, f"Expected 14 slides, found {len(prs.slides)}"
    for index, slide in enumerate(prs.slides, start=1):
        assert len(slide.shapes) > 0, f"Slide {index} is empty"
        for shape in slide.shapes:
            assert shape.left >= 0 and shape.top >= 0, f"Slide {index}: negative shape position"
            assert shape.left + shape.width <= prs.slide_width + Inches(0.02), f"Slide {index}: shape overflows right"
            assert shape.top + shape.height <= prs.slide_height + Inches(0.02), f"Slide {index}: shape overflows bottom"


if __name__ == "__main__":
    validate_assets()
    validate_deck()
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    # Re-open the final package to catch relationship or XML errors before handoff.
    reopened = Presentation(OUT)
    assert len(reopened.slides) == 14
    print(f"Created {OUT}")
    print(f"Slides: {len(reopened.slides)} | Size: {OUT.stat().st_size:,} bytes")
